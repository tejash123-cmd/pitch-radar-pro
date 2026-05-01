"""Extract text from diligence decks (PDF, PPTX) and optional founder website snippets."""
from __future__ import annotations

import io
import re
from urllib.parse import urlparse

import requests
from bs4 import BeautifulSoup
from pypdf import PdfReader

MAX_DECK_CHARS = 48_000
MAX_WEB_CHARS = 5_500
FETCH_TIMEOUT = 18


def _clean_ws(text: str) -> str:
    text = (text or "").replace("\x00", " ")
    return re.sub(r"[ \t]+", " ", re.sub(r"\n{3,}", "\n\n", text.strip()))


def extract_pdf_bytes(data: bytes) -> str:
    reader = PdfReader(io.BytesIO(data))
    parts: list[str] = []
    for page in reader.pages:
        chunk = (page.extract_text() or "").strip()
        if chunk:
            parts.append(chunk)
        if sum(len(x) for x in parts) > MAX_DECK_CHARS:
            break
    return _clean_ws("\n\n".join(parts))


def extract_pptx_bytes(data: bytes) -> str:
    try:
        from pptx import Presentation
    except ImportError as e:
        raise RuntimeError("python-pptx is required for .pptx files") from e

    prs = Presentation(io.BytesIO(data))
    blocks: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    t = "".join(run.text or "" for run in paragraph.runs).strip()
                    if t:
                        blocks.append(t)
            elif getattr(shape, "text", None):
                t = (shape.text or "").strip()
                if t:
                    blocks.append(t)
    return _clean_ws("\n".join(blocks))


def extract_deck_text(filename: str, data: bytes) -> tuple[str, str]:
    fn = (filename or "").strip().lower()
    if fn.endswith(".pdf"):
        txt = extract_pdf_bytes(data)
        return "PDF", txt
    if fn.endswith(".pptx"):
        txt = extract_pptx_bytes(data)
        return "PowerPoint", txt
    raise ValueError("Deck must be a .pdf or .pptx file")


def truncate_for_prompt(text: str, max_chars: int = MAX_DECK_CHARS) -> str:
    text = text or ""
    if len(text) <= max_chars:
        return text
    head_len = max_chars * 2 // 3
    tail_len = max_chars - head_len - 48
    return text[:head_len] + "\n\n[ … middle truncated for model context … ]\n\n" + text[-tail_len:]


def fetch_website_text(url_raw: str) -> tuple[str, str]:
    url_raw = (url_raw or "").strip()
    if not url_raw:
        return "", ""

    parsed = urlparse(url_raw if "://" in url_raw else f"https://{url_raw}")
    if parsed.scheme not in ("http", "https") or not parsed.netloc:
        return "", "Invalid URL (use https://…)"

    try:
        resp = requests.get(
            parsed.geturl(),
            timeout=FETCH_TIMEOUT,
            headers={
                "User-Agent": "TechnologyForesightAI/VC-diligence (contact: localhost)",
                "Accept": "text/html,application/xhtml+xml",
            },
        )
        resp.raise_for_status()
        ctype = (resp.headers.get("Content-Type") or "").lower()
        if "html" not in ctype and "text" not in ctype:
            return "", f"Non-HTML content ({ctype or 'unknown'})"
        soup = BeautifulSoup(resp.text[:800_000], "lxml")
        for tag in soup(["script", "style", "noscript"]):
            tag.decompose()
        text = soup.get_text("\n")
        text = _clean_ws(text)
        if len(text) > MAX_WEB_CHARS:
            text = text[: MAX_WEB_CHARS - 20] + "… [truncated]"
        return text, ""
    except Exception as e:
        return "", str(e)


def build_diligence_dossier(
    deck_filename: str,
    deck_inner_label: str,
    deck_raw_text: str,
    meeting_notes: str = "",
    website_url: str = "",
    website_text: str = "",
    extra_context: str = "",
    website_note: str = "",
) -> str:
    sections: list[str] = []

    headline = truncate_for_prompt(deck_raw_text, MAX_DECK_CHARS)
    sections.append(
        f"[STARTUP PITCH MATERIAL — uploaded {deck_inner_label}: {deck_filename}]\n{headline}"
    )

    notes = _clean_ws(meeting_notes)
    if notes:
        sections.append(f"[OPTIONAL FOUNDER / VC MEETING NOTES]\n{notes}")

    if website_url.strip():
        if website_text.strip():
            hdr = "[OPTIONAL COMPANY WEBSITE CONTENT]"
            hdr += f"\n(Source URL {website_url})"
            if website_note:
                hdr += f"\n(Fetch note: {website_note})"
            sections.append(f"{hdr}\n{website_text}")
        else:
            note = website_note or "No HTML text extracted."
            sections.append(
                "[OPTIONAL COMPANY WEBSITE — FETCH INCOMPLETE]\n"
                f"URL attempted: {website_url}\n{note}"
            )

    xc = _clean_ws(extra_context)
    if xc:
        sections.append(f"[OPTIONAL EXTRA CONTEXT FROM ANALYST]\n{xc}")

    return "\n\n---\n\n".join(sections)


def dossier_stub_for_store(dossier: str, deck_file: str) -> dict:
    return {
        "deck_file": deck_file,
        "dossier_char_count": len(dossier or ""),
        "dossier_preview": (dossier or "")[:800],
    }
